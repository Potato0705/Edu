from __future__ import annotations

import csv
from pathlib import Path

import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt


ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "reports"
ASSET_DIR = OUT_DIR / "assets"
OUT_PPT = OUT_DIR / "PACE_AES_项目汇报_20260427.pptx"


PROMPTS = [1, 2, 3, 4, 5, 6, 7, 8]
OLD_WISE = [0.393570, 0.531509, 0.330745, 0.534836, 0.500432, 0.536031, 0.492478, 0.549437]
LOCAL_WISE = [0.351742, 0.528236, 0.564445, 0.502703, 0.560138, 0.574416, 0.593228, 0.417740]
PACE_QWK = [0.793242, 0.603974, 0.612661, 0.711011, 0.743834, 0.717967, 0.739641, 0.549772]
PACE_MAE = [0.812325, 0.594444, 0.604046, 0.570621, 0.523546, 0.513889, 2.554140, 4.000000]
DELTA_QWK = [0.441500, 0.075738, 0.048216, 0.208308, 0.183695, 0.143551, 0.146413, 0.132032]
DELTA_QWK_OLD = [0.399672, 0.072466, 0.281916, 0.176174, 0.243402, 0.181936, 0.247162, 0.000335]


BG = RGBColor(247, 244, 236)
INK = RGBColor(34, 39, 46)
ACCENT = RGBColor(185, 68, 33)
ACCENT_2 = RGBColor(42, 90, 130)
ACCENT_3 = RGBColor(96, 122, 52)
MUTED = RGBColor(108, 110, 115)
GRID = RGBColor(222, 214, 202)


def ensure_dirs() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)


def build_summary_rows() -> list[dict[str, float | str]]:
    rows: list[dict[str, float | str]] = []
    for i, p in enumerate(PROMPTS):
        rows.append(
            {
                "Prompt": f"P{p}",
                "Old WISE QWK": OLD_WISE[i],
                "Local-WISE QWK": LOCAL_WISE[i],
                "PACE QWK": PACE_QWK[i],
                "ΔQWK vs Local": DELTA_QWK[i],
                "ΔQWK vs Old": DELTA_QWK_OLD[i],
                "PACE MAE": PACE_MAE[i],
            }
        )
    rows.append(
        {
            "Prompt": "平均",
            "Old WISE QWK": sum(OLD_WISE) / len(OLD_WISE),
            "Local-WISE QWK": sum(LOCAL_WISE) / len(LOCAL_WISE),
            "PACE QWK": sum(PACE_QWK) / len(PACE_QWK),
            "ΔQWK vs Local": sum(DELTA_QWK) / len(DELTA_QWK),
            "ΔQWK vs Old": sum(DELTA_QWK_OLD) / len(DELTA_QWK_OLD),
            "PACE MAE": sum(PACE_MAE) / len(PACE_MAE),
        }
    )
    return rows


def write_summary_csv(rows: list[dict[str, float | str]]) -> Path:
    path = OUT_DIR / "fold0_分数汇总表.csv"
    headers = list(rows[0].keys())
    with path.open("w", encoding="utf-8-sig", newline="") as f:
        writer = csv.DictWriter(f, fieldnames=headers)
        writer.writeheader()
        for row in rows:
            writer.writerow(row)
    return path


def make_summary_table_image(rows: list[dict[str, float | str]]) -> Path:
    fig, ax = plt.subplots(figsize=(12.8, 4.8), dpi=180)
    ax.axis("off")
    headers = list(rows[0].keys())
    cell_text = []
    for row in rows:
        rendered = []
        for h in headers:
            v = row[h]
            if isinstance(v, float):
                rendered.append(f"{v:.3f}")
            else:
                rendered.append(str(v))
        cell_text.append(rendered)
    tbl = ax.table(
        cellText=cell_text,
        colLabels=headers,
        loc="center",
        cellLoc="center",
        colLoc="center",
    )
    tbl.auto_set_font_size(False)
    tbl.set_fontsize(10)
    tbl.scale(1.05, 1.6)
    n_rows = len(rows)
    for (r, c), cell in tbl.get_celld().items():
        cell.set_edgecolor("#D8D0C6")
        if r == 0:
            cell.set_facecolor("#E8DED1")
            cell.set_text_props(weight="bold", color="#22272E")
        elif r == n_rows:
            cell.set_facecolor("#F4ECE2")
            cell.set_text_props(weight="bold", color="#B94421")
        elif c in [4, 5]:
            text = cell.get_text().get_text()
            if text.startswith("-"):
                cell.set_facecolor("#F4E8E3")
            else:
                cell.set_facecolor("#E9F1E1")
        else:
            cell.set_facecolor("#FFFFFF")
    ax.set_title("Fold0 分数汇总表（QWK 为主）", fontsize=16, fontweight="bold", pad=12)
    path = ASSET_DIR / "fold0_score_summary_table.png"
    fig.savefig(path, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    return path


def make_charts() -> dict[str, Path]:
    plt.rcParams["font.sans-serif"] = ["Microsoft YaHei", "SimHei", "Arial Unicode MS"]
    plt.rcParams["axes.unicode_minus"] = False

    charts: dict[str, Path] = {}

    # Chart 1: grouped QWK by prompt
    fig, ax = plt.subplots(figsize=(10.5, 4.8), dpi=180)
    x = range(len(PROMPTS))
    w = 0.24
    ax.bar([i - w for i in x], LOCAL_WISE, width=w, color="#6E7C8C", label="Local-WISE replay")
    ax.bar(x, PACE_QWK, width=w, color="#B94421", label="PACE (qwkfix)")
    ax.bar([i + w for i in x], OLD_WISE, width=w, color="#5F7A34", label="Old WISE ref")
    ax.set_xticks(list(x), [f"P{p}" for p in PROMPTS])
    ax.set_ylim(0, 0.9)
    ax.set_ylabel("QWK")
    ax.set_title("Fold0 各 Prompt 的 QWK 对比")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False, ncol=3, loc="upper center")
    fig.tight_layout()
    p = ASSET_DIR / "fold0_qwk_grouped.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    charts["qwk_grouped"] = p

    # Chart 2: delta qwk vs local
    fig, ax = plt.subplots(figsize=(10.5, 4.2), dpi=180)
    colors = ["#B94421" if v >= 0 else "#7A7A7A" for v in DELTA_QWK]
    ax.bar([f"P{p}" for p in PROMPTS], DELTA_QWK, color=colors)
    ax.axhline(0, color="#333333", linewidth=1)
    ax.set_ylabel("ΔQWK vs Local-WISE")
    ax.set_title("修复后 Fold0：PACE 相对 Local-WISE 的 QWK 提升")
    ax.grid(axis="y", alpha=0.25)
    for i, v in enumerate(DELTA_QWK):
        ax.text(i, v + 0.01, f"{v:+.3f}", ha="center", va="bottom", fontsize=8)
    fig.tight_layout()
    p = ASSET_DIR / "fold0_delta_qwk.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    charts["delta_qwk"] = p

    # Chart 3: MAE compare
    fig, ax = plt.subplots(figsize=(10.5, 4.4), dpi=180)
    ax.plot([f"P{p}" for p in PROMPTS], [2.061625, 0.547222, 0.479769, 0.522599, 0.506925, 0.575000, 3.254777, 6.268966],
            marker="o", linewidth=2.2, color="#6E7C8C", label="Local-WISE MAE")
    ax.plot([f"P{p}" for p in PROMPTS], PACE_MAE,
            marker="o", linewidth=2.2, color="#B94421", label="PACE MAE")
    ax.set_ylabel("MAE")
    ax.set_title("Fold0 各 Prompt 的 MAE 对比")
    ax.grid(axis="y", alpha=0.25)
    ax.legend(frameon=False)
    fig.tight_layout()
    p = ASSET_DIR / "fold0_mae_lines.png"
    fig.savefig(p, bbox_inches="tight", facecolor="white")
    plt.close(fig)
    charts["mae_lines"] = p

    return charts


def set_slide_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def style_run(run, size: int = 20, bold: bool = False, color: RGBColor = INK, font_name: str = "Microsoft YaHei"):
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    tb = slide.shapes.add_textbox(Cm(1.2), Cm(0.7), Cm(30.5), Cm(1.8))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    style_run(r, size=24, bold=True, color=INK)
    if subtitle:
        st = slide.shapes.add_textbox(Cm(1.2), Cm(2.0), Cm(30.0), Cm(0.9))
        tf2 = st.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        style_run(r2, size=10, color=MUTED)


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, level_sizes: tuple[int, int] = (18, 15)) -> None:
    tb = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(6)
        r = p.add_run()
        r.text = item
        style_run(r, size=level_sizes[0], color=INK)


def add_note_box(slide, title: str, body: list[str], left: float, top: float, width: float, height: float, accent: RGBColor = ACCENT_2) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(left), Cm(top), Cm(width), Cm(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.clear()
    p0 = tf.paragraphs[0]
    r0 = p0.add_run()
    r0.text = title
    style_run(r0, size=16, bold=True, color=accent)
    for text in body:
        p = tf.add_paragraph()
        p.level = 0
        p.space_before = Pt(4)
        r = p.add_run()
        r.text = text
        style_run(r, size=13, color=INK)


def add_flow_box(slide, x, y, w, h, title, body, fill_rgb, line_rgb=INK):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Cm(x), Cm(y), Cm(w), Cm(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.color.rgb = line_rgb
    shape.line.width = Pt(1.2)
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = title
    style_run(r, size=16, bold=True, color=INK)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = body
    style_run(r2, size=12, color=INK)
    return shape


def add_connector(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    line.line.color.rgb = color
    line.line.width = Pt(2.0)


def build_presentation(charts: dict[str, Path]) -> Presentation:
    prs = Presentation()
    prs.slide_width = Cm(33.867)
    prs.slide_height = Cm(19.05)
    summary_rows = build_summary_rows()

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(19.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(250, 247, 240)
    bar.line.fill.background()
    accent_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Cm(0), Cm(0), Cm(33.867), Cm(1.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = ACCENT
    accent_bar.line.fill.background()
    tb = slide.shapes.add_textbox(Cm(1.3), Cm(3.0), Cm(29.5), Cm(4.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "PACE-AES 项目阶段汇报"
    style_run(r, size=30, bold=True, color=INK)
    p2 = tb.text_frame.add_paragraph()
    r2 = p2.add_run()
    r2.text = "围绕 WISE-AES 的二层校准：动机、方法、初步结果与下一步"
    style_run(r2, size=18, color=ACCENT_2)
    tb2 = slide.shapes.add_textbox(Cm(1.3), Cm(14.8), Cm(20), Cm(1.6))
    p = tb2.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "汇报时间：2026-04-27    汇报对象：师兄"
    style_run(r, size=14, color=MUTED)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "1. 立项动机", "为什么要在 WISE-AES 之上继续做 PACE-AES")
    add_bullets(
        slide,
        [
            "WISE-AES 已经解决了 protocol alignment：通过 rubric induction + exemplar evolution，把 LLM 评分协议调到可用水平。",
            "但 WISE-AES 仍有明显残余误差，尤其集中在分数边界附近：同一篇作文经常“方向对了，但落点不稳”。",
            "我们的核心判断：protocol alignment 不等于 score alignment，Layer 1 之后仍需要一个专门的分数校准层。",
            "因此提出 PACE-AES：固定 WISE 的 evolved protocol，只做 Layer 2 的 protocol-conditioned ordinal calibration。",
        ],
        left=1.4,
        top=3.1,
        width=16.2,
        height=11.5,
    )
    add_note_box(
        slide,
        "本项目想解决的不是",
        [
            "不是重新证明 WISE-AES 本身有效。",
            "不是重新跑一遍完整 evolution。",
            "而是回答：在固定 WISE protocol 后，剩余误差能否被系统性修正？",
        ],
        left=19.0,
        top=4.0,
        width=13.0,
        height=8.0,
        accent=ACCENT,
    )

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "2. 当前方法框架", "WISE-AES 作为 Layer 1，PACE-AES 作为 Layer 2/3")
    add_flow_box(slide, 1.3, 5.0, 7.0, 4.0, "Layer 1\nWISE-AES", "evolved rubric\n+ static anchors", RGBColor(235, 227, 215))
    add_flow_box(slide, 10.0, 5.0, 8.2, 4.0, "Local-WISE Replay", "同一 protocol 下重新打分\n得到 y_raw 与 hidden state", RGBColor(227, 236, 243))
    add_flow_box(slide, 20.0, 5.0, 5.8, 4.0, "Evidence", "anchor-relative\nlatent evidence", RGBColor(231, 240, 225))
    add_flow_box(slide, 27.0, 5.0, 5.0, 4.0, "PACE", "ordinal calibrator\n输出 y_pred", RGBColor(248, 229, 221))
    add_connector(slide, 8.3, 7.0, 10.0, 7.0)
    add_connector(slide, 18.2, 7.0, 20.0, 7.0)
    add_connector(slide, 25.8, 7.0, 27.0, 7.0)
    add_note_box(
        slide,
        "核心设计点",
        [
            "输入不是只有 h(x)，而是 h(x) 相对 low/mid/high anchors 的残差坐标。",
            "输出不是普通回归，而是显式有序边界建模。",
            "最终目标是提升 QWK，同时尽量不破坏分数结构。",
        ],
        left=2.0,
        top=11.3,
        width=29.0,
        height=4.8,
        accent=ACCENT_3,
    )

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "3. 目前实现到哪一步", "代码、数据与实验基础设施")
    add_bullets(
        slide,
        [
            "已完成本地 / 服务器双端 PACE 代码：`pace/evidence.py`、`pace/calibration.py`、`train_pace_calibrator.py`、`sweep_pace_calibrator.py`。",
            "已上传并复用完整 WISE Layer-1 logs：ASAP 全部 8 个 prompt × 5 folds，共 40 个 runs。",
            "已实现 anchor cache、evidence cache、split 级断点续跑、sweep 汇总、Prompt/Fold 定位等基础设施。",
            "当前 AutoDL 服务器使用 RTX 5090 32GB，本地 Llama-3.1-8B-Instruct 作为统一后端，不再依赖 OpenRouter API。",
        ],
        left=1.4,
        top=3.2,
        width=17.3,
        height=10.5,
    )
    add_note_box(
        slide,
        "实验口径",
        [
            "严格主比较：Local-WISE replay vs PACE-AES",
            "参考比较：Original WISE-AES old logs",
            "这样可以避免把不同后端输出直接混成一个 baseline。",
        ],
        left=19.5,
        top=4.1,
        width=11.8,
        height=5.2,
        accent=ACCENT_2,
    )
    add_note_box(
        slide,
        "当前已知风险",
        [
            "不同 prompt 的最优配置不完全一致。",
            "小范围 prompt 与大范围 prompt 可能需要不同 decode 策略。",
        ],
        left=19.5,
        top=10.0,
        width=11.8,
        height=4.4,
        accent=ACCENT,
    )

    # Slide 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "4. 关键问题定位：QWK loss 实现错误", "这一步是当前阶段最重要的工程修复")
    add_note_box(
        slide,
        "发现的问题",
        [
            "原先的 soft_qwk_loss 实际写成了 soft_qwk_score 本身。",
            "训练时等价于在最小化 QWK，而不是最大化 QWK。",
            "同时又被 clamp 到 0，导致日志里大量出现 loss_qwk = 0.0。",
        ],
        left=1.4,
        top=3.3,
        width=14.6,
        height=6.3,
        accent=ACCENT,
    )
    add_note_box(
        slide,
        "修复内容",
        [
            "改为：soft_qwk_score = 1 - num / den",
            "soft_qwk_loss = 1 - soft_qwk_score = num / den",
            "并把 soft_qwk_score 直接写进 summary，便于排查训练是否真正生效。",
        ],
        left=17.0,
        top=3.3,
        width=14.5,
        height=6.3,
        accent=ACCENT_2,
    )
    add_note_box(
        slide,
        "修复后的直接效果",
        [
            "P7 / P8 这种大分数范围 prompt 不再出现 QWK 崩盘。",
            "QWK 项开始真正参与训练，而不是只靠 ordinal CE 在“误打误撞”。",
        ],
        left=6.2,
        top=11.0,
        width=20.5,
        height=4.2,
        accent=ACCENT_3,
    )

    # Slide 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "5. Fold0 初步结果：QWK 全部正增益", "修复后的诊断汇总")
    slide.shapes.add_picture(str(charts["qwk_grouped"]), Cm(1.2), Cm(3.0), width=Cm(31.0))
    tb = slide.shapes.add_textbox(Cm(1.6), Cm(14.5), Cm(30.0), Cm(2.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "结论：相对 Local-WISE replay，Fold0 上 8/8 个 prompt 的 QWK 全部提升。"
    style_run(r, size=16, bold=True, color=ACCENT)

    # Slide 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "6. Fold0 汇总：提升幅度", "均值已经具备论文价值")
    slide.shapes.add_picture(str(charts["delta_qwk"]), Cm(1.2), Cm(3.0), width=Cm(31.0))
    add_note_box(
        slide,
        "当前诊断均值",
        [
            "mean ΔQWK vs Local-WISE = +0.172",
            "mean ΔQWK vs old WISE ref = +0.200",
            "P1/P4/P5/P7/P8 的提升尤为明显。",
        ],
        left=2.0,
        top=12.8,
        width=13.5,
        height=3.5,
        accent=ACCENT,
    )
    add_note_box(
        slide,
        "解释口径",
        [
            "这张表是很强的 fold0 diagnostic summary。",
            "最终主表仍需补齐剩余 folds，避免只在单折上结论过强。",
        ],
        left=17.0,
        top=12.8,
        width=14.0,
        height=3.5,
        accent=ACCENT_2,
    )

    # Slide 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "7. 两类 Prompt 的策略差异", "这一点决定后续 5-fold 怎么跑")
    add_note_box(
        slide,
        "小/中范围 Prompt（P1-P6）",
        [
            "目前更适合 threshold decode。",
            "典型配置：lr 1e-3, lambda_qwk 0.25。",
            "P2 例外：lr 1e-4 更稳，但依然不需要 blend。",
        ],
        left=1.5,
        top=4.0,
        width=13.5,
        height=7.5,
        accent=ACCENT_2,
    )
    add_note_box(
        slide,
        "大范围 Prompt（P7-P8）",
        [
            "需要 QWK loss 修复 + 保守 blend decode。",
            "当前有效配置：lr 3e-4, lambda_qwk 2.0, blend_alpha 0.65, max_delta_frac 0.10。",
            "P8 已从 0.230 崩盘修到 0.550，P7 修到 0.740。",
        ],
        left=17.2,
        top=4.0,
        width=14.0,
        height=7.5,
        accent=ACCENT,
    )
    add_note_box(
        slide,
        "启发",
        [
            "PACE 并不是单一配置走天下。",
            "更像是“同一方法骨架 + prompt-wise 稳定化配置”。",
        ],
        left=6.8,
        top=12.4,
        width=19.0,
        height=3.0,
        accent=ACCENT_3,
    )

    # Slide 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "8. 当前进度与下一步", "截至 2026-04-27")
    add_bullets(
        slide,
        [
            "已完成：Fold0 全部 prompt 的诊断结果；QWK loss bug 修复；P2/P7/P8 的低成本结构增强验证。",
            "正在进行：P7/P8 的 fold1-4 5-fold 实验，用于先把大分数范围 prompt 的故事做扎实。",
            "下一步候选 A：补齐 P1-P6 的剩余 folds，形成完整 5-fold 主表。",
            "下一步候选 B：先围绕 P7/P8 写一个“wide-range prompt 校准收益”子故事，再回头补小范围 prompt。",
        ],
        left=1.4,
        top=3.2,
        width=18.0,
        height=10.8,
    )
    add_note_box(
        slide,
        "建议决策",
        [
            "先把 P7/P8 的 fold1-4 跑完。",
            "如果 5-fold 稳定，再决定是否优先扩到全部 prompts。",
            "这能最快形成一条可写进论文的稳定主线。",
        ],
        left=20.0,
        top=4.3,
        width=11.5,
        height=6.0,
        accent=ACCENT_3,
    )

    # Slide 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "9. 与师兄讨论的三个问题", "建议本次汇报重点确认")
    add_bullets(
        slide,
        [
            "论文主线是否先围绕“WISE protocol 之上的二层校准”收紧，而不同时展开过多增强项？",
            "实验资源分配上，是否优先把 P7/P8 的 5-fold 跑完整，再决定是否扩展到全部 prompts？",
            "最终论文口径上，是否接受“Local-WISE replay 作为严格主 baseline，old WISE log 作为 reference baseline”？",
        ],
        left=1.6,
        top=4.0,
        width=29.5,
        height=8.5,
    )
    tb = slide.shapes.add_textbox(Cm(1.6), Cm(14.4), Cm(30.0), Cm(1.2))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "一句话结论：方法已经站住，当前重点从“能不能跑通”转向“如何把 5-fold 证据链补齐”。"
    style_run(r, size=17, bold=True, color=ACCENT)

    # Appendix slide: score summary table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide)
    add_title(slide, "附录：Fold0 分数汇总表", "便于和师兄快速对齐当前量化结果")
    table_img = make_summary_table_image(summary_rows)
    slide.shapes.add_picture(str(table_img), Cm(1.0), Cm(2.8), width=Cm(31.7))
    tb = slide.shapes.add_textbox(Cm(1.2), Cm(16.2), Cm(31.0), Cm(1.0))
    p = tb.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = "说明：此表基于当前 fold0 诊断结果，P7/P8 使用 qwkfix 版本，后续仍需补齐剩余 folds。"
    style_run(r, size=12, color=MUTED)

    return prs


def main() -> None:
    ensure_dirs()
    summary_rows = build_summary_rows()
    write_summary_csv(summary_rows)
    charts = make_charts()
    prs = build_presentation(charts)
    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
